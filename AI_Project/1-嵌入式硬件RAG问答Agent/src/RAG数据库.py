"""
嵌入式RAG问答助手 - 向量数据库构建模块
功能：解析硬件文档并构建FAISS向量库
使用DashScope API进行文本向量化
"""

import os
import sys
import json
import pickle
from pathlib import Path
from typing import List, Dict, Any
import numpy as np
import time

# 设置标准输出编码为UTF-8
if sys.platform == 'win32':
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

# 文档解析库
import PyPDF2
from pptx import Presentation
import openpyxl
from docx import Document  # 添加DOCX支持
import hashlib

# 向量化和存储
import dashscope
from dashscope import TextEmbedding
import faiss

# 混合检索支持
try:
    from rank_bm25 import BM25Okapi
    import jieba
    BM25_AVAILABLE = True
except ImportError:
    BM25_AVAILABLE = False
    print("⚠️  BM25库未安装，将只使用语义检索。安装: pip install rank-bm25 jieba")

# 缓存支持
try:
    # 添加项目根目录到路径
    current_dir = os.path.dirname(os.path.abspath(__file__))
    project_root = os.path.dirname(current_dir)
    if project_root not in sys.path:
        sys.path.insert(0, project_root)
    
    from utils.cache_manager import get_embedding_cache
    CACHE_AVAILABLE = True
except ImportError:
    CACHE_AVAILABLE = False
    print("⚠️  缓存模块未找到，将不使用 Embedding 缓存")


class DocumentParser:
    """文档解析器：支持PDF、PPTX、XLSX"""
    
    @staticmethod
    def parse_pdf(file_path: str) -> str:
        """解析PDF文件"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        except Exception as e:
            print(f"❌ 解析PDF失败 [{file_path}]: {e}")
        return text
    
    @staticmethod
    def parse_pptx(file_path: str) -> str:
        """解析PPTX文件"""
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"❌ 解析PPTX失败 [{file_path}]: {e}")
        return text
    
    @staticmethod
    def parse_xlsx(file_path: str) -> str:
        """解析XLSX文件"""
        text = ""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
        except Exception as e:
            print(f"❌ 解析XLSX失败 [{file_path}]: {e}")
        return text
    
    @staticmethod
    def parse_docx(file_path: str) -> str:
        """解析DOCX文件"""
        text = ""
        try:
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            # 也提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    if row_text.strip():
                        text += row_text + "\n"
        except Exception as e:
            print(f"❌ 解析DOCX失败 [{file_path}]: {e}")
        return text
    
    @classmethod
    def parse_document(cls, file_path: str) -> str:
        """根据文件类型自动选择解析方法"""
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pdf':
            return cls.parse_pdf(file_path)
        elif ext == '.pptx':
            return cls.parse_pptx(file_path)
        elif ext in ['.xlsx', '.xls']:
            return cls.parse_xlsx(file_path)
        elif ext == '.docx':
            return cls.parse_docx(file_path)
        else:
            print(f"⚠️  不支持的文件类型: {ext}")
            return ""


class TextChunker:
    """文本分块器：将长文本切分成合适的块"""
    
    @staticmethod
    def tokenize(text: str) -> List[str]:
        """
        中文分词
        
        Args:
            text: 原始文本
            
        Returns:
            分词后的词列表
        """
        if not BM25_AVAILABLE:
            # 如果jieba不可用，简单按字符分割
            return list(text)
        
        # 使用jieba分词
        tokens = jieba.lcut(text)
        # 过滤停用词和标点
        tokens = [t.strip() for t in tokens if t.strip() and len(t.strip()) > 1]
        return tokens
    
    @staticmethod
    def clean_text(text: str) -> str:
        """
        安全地清理文本，避免内存问题
        
        Args:
            text: 原始文本
            
        Returns:
            清理后的文本
        """
        import re
        
        # 替换多个连续空白为单个空格
        text = re.sub(r'\s+', ' ', text)
        # 去除首尾空白
        text = text.strip()
        return text
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
        """
        智能文本分块（内存优化版）
        
        Args:
            text: 原始文本
            chunk_size: 每块的字符数
            overlap: 块之间的重叠字符数
            
        Returns:
            文本块列表
        """
        if not text or not text.strip():
            return []
        
        # 安全地清理文本
        try:
            text = TextChunker.clean_text(text)
        except Exception as e:
            print(f"  ⚠️  文本清理失败: {e}，使用原始文本")
        
        text_length = len(text)
        
        # 如果文本很短，直接返回
        if text_length <= chunk_size:
            return [text] if text else []
        
        chunks = []
        start = 0
        max_chunks = (text_length // chunk_size) * 2 + 10  # 预防无限循环
        chunk_count = 0
        
        while start < text_length and chunk_count < max_chunks:
            end = min(start + chunk_size, text_length)
            
            # 如果不是最后一块，尝试在句子边界处切分
            if end < text_length:
                # 寻找最近的句子结束标记
                best_end = end
                for delimiter in ['. ', '。 ', '! ', '！ ', '? ', '？ ', '\n', '。', '.']:
                    pos = text.rfind(delimiter, start, end)
                    if pos != -1 and pos > start:
                        best_end = pos + len(delimiter)
                        break
                end = best_end
            
            # 提取块
            try:
                chunk = text[start:end].strip()
                if chunk and len(chunk) > 10:  # 过滤太短的块
                    chunks.append(chunk)
            except Exception as e:
                print(f"  ⚠️  提取文本块失败: {e}")
                break
            
            # 计算下一个起始位置（考虑重叠）
            start = max(start + 1, end - overlap) if end < text_length else text_length
            chunk_count += 1
        
        return chunks


class EmbeddingVectorDatabase:
    """向量数据库：使用FAISS进行高效向量检索，使用DashScope API进行向量化"""
    
    def __init__(self, 
                 model_name: str = "text-embedding-v2",
                 vector_db_path: str = "./embedding向量库"):
        """
        初始化向量数据库
        
        Args:
            model_name: DashScope Embedding模型名称
            vector_db_path: 向量库保存路径
        """
        print("🚀 初始化嵌入式RAG向量数据库...")
        
        # 从环境变量获取API Key
        api_key = os.getenv("DASHSCOPE_API_KEY")
        if not api_key:
            raise ValueError("❌ 未找到环境变量 DASHSCOPE_API_KEY，请先设置")
        
        dashscope.api_key = api_key
        print("✅ DashScope API Key 已加载")
        
        self.model_name = model_name
        self.vector_db_path = Path(vector_db_path)
        self.vector_db_path.mkdir(exist_ok=True)
        
        # DashScope text-embedding-v2 模型的向量维度是1536
        self.embedding_dim = 1536
        print(f"📊 使用模型: {model_name}, 向量维度: {self.embedding_dim}")
        
        # 初始化FAISS索引
        self.index = None
        self.documents = []  # 存储文档元数据
        self.processed_files = {}  # 记录已处理文件: {filename: {hash, mtime, chunk_count}}
        
        # BM25索引（用于关键词检索）
        self.bm25 = None
        self.tokenized_docs = []  # 分词后的文档，用于BM25
        
        # 初始化 Embedding 缓存
        self.embedding_cache = None
        if CACHE_AVAILABLE:
            try:
                self.embedding_cache = get_embedding_cache()
                print("✅ Embedding 缓存已启用")
            except Exception as e:
                print(f"⚠️  Embedding 缓存初始化失败: {e}")
    
    def embed_texts(self, texts: List[str], batch_size: int = 25) -> np.ndarray:
        """
        使用DashScope API批量向量化文本（支持缓存）
        
        Args:
            texts: 文本列表
            batch_size: 每批处理的文本数量（DashScope限制每次最多25条）
            
        Returns:
            向量矩阵 (n_texts, embedding_dim)
        """
        all_embeddings = []
        texts_to_embed = []  # 需要调用 API 的文本
        text_indices = []    # 记录需要调用 API 的文本索引
        
        # 1. 检查缓存
        cache_hits = 0
        if self.embedding_cache:
            cached_results = self.embedding_cache.get_batch(texts)
            for idx, text in enumerate(texts):
                cached_emb = cached_results.get(text)
                if cached_emb is not None:
                    all_embeddings.append(cached_emb)
                    cache_hits += 1
                else:
                    all_embeddings.append(None)  # 占位符
                    texts_to_embed.append(text)
                    text_indices.append(idx)
        else:
            # 没有缓存，全部需要调用 API
            texts_to_embed = texts
            text_indices = list(range(len(texts)))
            all_embeddings = [None] * len(texts)
        
        # 显示缓存命中情况
        if self.embedding_cache:
            cache_miss = len(texts_to_embed)
            print(f"📊 向量化 {len(texts)} 个文本块 (缓存命中: {cache_hits}, 需要API调用: {cache_miss})")
        else:
            print(f"📊 开始向量化 {len(texts)} 个文本块...")
        
        # 2. 批量调用 API 处理未缓存的文本
        if texts_to_embed:
            total_batches = (len(texts_to_embed) + batch_size - 1) // batch_size
            new_embeddings = []
            
            for i in range(0, len(texts_to_embed), batch_size):
                batch = texts_to_embed[i:i + batch_size]
                batch_num = i // batch_size + 1
                
                try:
                    # 调用DashScope API
                    response = TextEmbedding.call(
                        model=self.model_name,
                        input=batch
                    )
                    
                    if response.status_code == 200:
                        # 提取embeddings
                        embeddings = [item['embedding'] for item in response.output['embeddings']]
                        new_embeddings.extend(embeddings)
                        print(f"  ✓ 批次 [{batch_num}/{total_batches}] 完成")
                    else:
                        print(f"  ✗ 批次 [{batch_num}/{total_batches}] 失败: {response.message}")
                        # 失败时用零向量填充
                        new_embeddings.extend([np.zeros(self.embedding_dim).tolist()] * len(batch))
                    
                    # API调用限流，避免过快
                    time.sleep(0.1)
                    
                except Exception as e:
                    print(f"  ✗ 批次 [{batch_num}/{total_batches}] 异常: {e}")
                    new_embeddings.extend([np.zeros(self.embedding_dim).tolist()] * len(batch))
            
            # 3. 将新获取的 embeddings 填充到结果中并保存到缓存
            for idx, emb_idx in enumerate(text_indices):
                all_embeddings[emb_idx] = new_embeddings[idx]
                
                # 保存到缓存
                if self.embedding_cache:
                    try:
                        self.embedding_cache.put(texts_to_embed[idx], new_embeddings[idx])
                    except Exception as e:
                        print(f"  ⚠️  缓存保存失败: {e}")
        
        return np.array(all_embeddings, dtype='float32')
    
    @staticmethod
    def get_short_path_name(long_path: str) -> str:
        """
        获取Windows短路径名（8.3格式），解决中文路径问题
        
        Args:
            long_path: 长路径名（可能包含中文）
            
        Returns:
            短路径名（纯ASCII字符）
        """
        if sys.platform != 'win32':
            return long_path
        
        try:
            import ctypes
            from ctypes import wintypes
            
            # 获取短路径名的Windows API
            _GetShortPathNameW = ctypes.windll.kernel32.GetShortPathNameW
            _GetShortPathNameW.argtypes = [wintypes.LPCWSTR, wintypes.LPWSTR, wintypes.DWORD]
            _GetShortPathNameW.restype = wintypes.DWORD
            
            # 确保路径存在
            if not os.path.exists(long_path):
                return long_path
            
            # 获取所需缓冲区大小
            buffer_size = _GetShortPathNameW(long_path, None, 0)
            if buffer_size == 0:
                return long_path
            
            # 获取短路径名
            short_path = ctypes.create_unicode_buffer(buffer_size)
            _GetShortPathNameW(long_path, short_path, buffer_size)
            return short_path.value
        except Exception:
            return long_path
    
    @staticmethod
    def get_file_hash(file_path: str) -> str:
        """计算文件MD5哈希值（快速版本，只读取前1MB）"""
        try:
            hash_md5 = hashlib.md5()
            with open(file_path, "rb") as f:
                # 只读取前1MB用于快速hash
                chunk = f.read(1024 * 1024)
                hash_md5.update(chunk)
            return hash_md5.hexdigest()
        except Exception:
            return ""
    
    def is_file_processed(self, file_path: Path) -> bool:
        """检查文件是否已处理过（通过hash和修改时间）"""
        filename = file_path.name
        if filename not in self.processed_files:
            return False
        
        # 检查文件修改时间
        current_mtime = file_path.stat().st_mtime
        stored_mtime = self.processed_files[filename].get('mtime', 0)
        
        if current_mtime != stored_mtime:
            return False  # 文件已修改
        
        # 检查文件hash
        current_hash = self.get_file_hash(str(file_path))
        stored_hash = self.processed_files[filename].get('hash', '')
        
        return current_hash == stored_hash
        
    def build_database(self, 
                      data_folder: str = "./data",
                      chunk_size: int = 500,
                      overlap: int = 50,
                      incremental: bool = True):
        """
        构建向量数据库（支持增量构建）
        
        Args:
            data_folder: 文档文件夹路径
            chunk_size: 文本块大小
            overlap: 块之间重叠大小
            incremental: 是否增量构建（True=只处理新文件，False=完全重建）
        
        Returns:
            是否构建成功
        """
        print("\n" + "="*60)
        if incremental:
            print("📚 开始增量构建向量数据库")
        else:
            print("📚 开始完全重建向量数据库")
        print("="*60)
        
        data_path = Path(data_folder)
        if not data_path.exists():
            print(f"❌ 数据文件夹不存在: {data_folder}")
            return False
        
        # 支持的文件类型（新增DOCX）
        supported_extensions = ['.pdf', '.pptx', '.xlsx', '.xls', '.docx']
        all_files = [f for f in data_path.iterdir() 
                     if f.is_file() and f.suffix.lower() in supported_extensions]
        
        if not all_files:
            print(f"⚠️  未找到支持的文档文件")
            return False
        
        # 筛选需要处理的文件
        if incremental:
            files = [f for f in all_files if not self.is_file_processed(f)]
            skipped_count = len(all_files) - len(files)
            print(f"📁 发现 {len(all_files)} 个文档文件")
            if skipped_count > 0:
                print(f"⏭️  跳过 {skipped_count} 个已处理文件")
            print(f"🆕 需要处理 {len(files)} 个新文件")
        else:
            files = all_files
            # 完全重建时，清空已处理文件记录
            self.processed_files = {}
            print(f"📁 发现 {len(files)} 个文档文件")
        
        if not files:
            if incremental:
                print(f"✅ 所有文件都已处理，无需更新")
                return True  # 增量构建时，没有新文件也返回True
            else:
                print(f"❌ 未找到支持的文档文件")
                return False
        
        all_chunks = []
        all_metadata = []
        
        # 遍历所有文档
        for idx, file_path in enumerate(files, 1):
            print(f"\n[{idx}/{len(files)}] 处理文档: {file_path.name}")
            
            # 解析文档
            print("  ├─ 解析文档内容...")
            text = DocumentParser.parse_document(str(file_path))
            
            if not text.strip():
                print("  └─ ⚠️  文档内容为空，跳过")
                continue
            
            print(f"  ├─ 提取文本: {len(text)} 字符")
            
            # 文本分块
            print(f"  ├─ 文本分块 (chunk_size={chunk_size}, overlap={overlap})...")
            chunks = TextChunker.chunk_text(text, chunk_size, overlap)
            print(f"  └─ ✅ 生成 {len(chunks)} 个文本块")
            
            # 保存块和元数据
            for chunk_idx, chunk in enumerate(chunks):
                all_chunks.append(chunk)
                all_metadata.append({
                    'filename': file_path.name,
                    'chunk_id': chunk_idx,
                    'text': chunk
                })
            
            # 记录文件为已处理
            self.processed_files[file_path.name] = {
                'hash': self.get_file_hash(str(file_path)),
                'mtime': file_path.stat().st_mtime,
                'chunk_count': len(chunks)
            }
        
        if not all_chunks:
            print("\n❌ 没有成功提取任何文本块")
            return False
        
        print(f"\n{'='*60}")
        print(f"📊 文本处理完成")
        print(f"  • 总文档数: {len(files)}")
        print(f"  • 总文本块: {len(all_chunks)}")
        print(f"{'='*60}")
        
        # 向量化新文本块
        if all_chunks:
            print("\n🔄 使用DashScope API向量化文本块...")
            new_embeddings = self.embed_texts(all_chunks, batch_size=25)
            print(f"✅ 向量化完成，向量形状: {new_embeddings.shape}")
            
            # 增量添加或重建FAISS索引
            if incremental and self.index is not None:
                print("\n🔨 增量更新FAISS索引...")
                self.index.add(new_embeddings.astype('float32'))
                self.documents.extend(all_metadata)
                print(f"✅ FAISS索引更新完成，现共 {self.index.ntotal} 个向量")
            else:
                print("\n🔨 构建FAISS索引...")
                self.index = faiss.IndexFlatL2(self.embedding_dim)
                self.index.add(new_embeddings.astype('float32'))
                self.documents = all_metadata
                print(f"✅ FAISS索引构建完成，共 {self.index.ntotal} 个向量")
            
            # 构建或更新BM25索引（用于混合检索）
            if BM25_AVAILABLE:
                print("\n🔨 构建BM25索引（关键词检索）...")
                if incremental and self.tokenized_docs:
                    # 增量模式：追加新的分词文档
                    new_tokenized = [TextChunker.tokenize(chunk) for chunk in all_chunks]
                    self.tokenized_docs.extend(new_tokenized)
                    self.bm25 = BM25Okapi(self.tokenized_docs)
                    print(f"✅ BM25索引更新完成，共 {len(self.tokenized_docs)} 个文档")
                else:
                    # 完全重建模式
                    self.tokenized_docs = [TextChunker.tokenize(chunk) for chunk in all_chunks]
                    self.bm25 = BM25Okapi(self.tokenized_docs)
                    print(f"✅ BM25索引构建完成，共 {len(self.tokenized_docs)} 个文档")
        else:
            print("\n⏭️  没有新文件需要处理，跳过向量化")
        
        # 保存向量库
        self.save_database()
        
        print(f"\n{'='*60}")
        print("🎉 向量数据库构建完成！")
        print(f"{'='*60}")
        
        return True
        
    def save_database(self):
        """保存向量数据库到磁盘（优化中文路径处理）"""
        print(f"\n💾 保存向量数据库到: {self.vector_db_path.absolute()}")
        
        # 确保目录存在
        self.vector_db_path.mkdir(parents=True, exist_ok=True)
        
        # 保存FAISS索引（多方法处理中文路径问题）
        index_path = self.vector_db_path / "faiss.index"
        faiss_saved = False
        
        # 方法1：使用Windows短路径名（8.3格式）
        if sys.platform == 'win32':
            try:
                abs_index_path = str(index_path.absolute())
                short_path = self.get_short_path_name(abs_index_path)
                
                if short_path != abs_index_path:
                    print(f"  ├─ 使用Windows短路径名保存...")
                    faiss.write_index(self.index, short_path)
                    if index_path.exists():
                        file_size = index_path.stat().st_size
                        print(f"  ├─ ✅ FAISS索引已保存（短路径）: {index_path.name} ({file_size:,} 字节)")
                        faiss_saved = True
            except Exception as e:
                print(f"  ├─ ⚠️  短路径方法失败: {e.__class__.__name__}")
        
        # 方法2：直接保存（如果短路径未成功）
        if not faiss_saved:
            try:
                abs_index_path = str(index_path.absolute())
                print(f"  ├─ 尝试直接保存FAISS索引...")
                faiss.write_index(self.index, abs_index_path)
                if index_path.exists():
                    file_size = index_path.stat().st_size
                    print(f"  ├─ ✅ FAISS索引已保存: {index_path.name} ({file_size:,} 字节)")
                    faiss_saved = True
            except Exception as e:
                print(f"  ├─ ⚠️  直接保存失败: {e.__class__.__name__}")
        
        # 方法3：通过临时文件（最后的备用方法）
        if not faiss_saved:
            try:
                print(f"  ├─ 使用临时文件备用方法...")
                import tempfile
                import shutil
                
                # 创建临时文件（系统临时目录，英文路径）
                with tempfile.NamedTemporaryFile(mode='wb', delete=False, suffix='.index') as tmp:
                    temp_path = tmp.name
                
                # 保存到临时文件
                faiss.write_index(self.index, temp_path)
                
                # 复制到目标位置
                abs_index_path = str(index_path.absolute())
                shutil.copy2(temp_path, abs_index_path)
                
                # 删除临时文件
                os.unlink(temp_path)
                
                # 验证
                if index_path.exists():
                    file_size = index_path.stat().st_size
                    print(f"  ├─ ✅ FAISS索引已保存（备用方法）: {index_path.name} ({file_size:,} 字节)")
                    faiss_saved = True
                else:
                    print(f"  ├─ ❌ 备用方法失败：文件未找到")
            except Exception as e:
                print(f"  ├─ ❌ 备用方法失败: {e}")
                import traceback
                traceback.print_exc()
        
        if not faiss_saved:
            print(f"  ├─ ❌ FAISS索引保存失败！所有方法均失败")
        
        # 保存文档元数据
        metadata_path = self.vector_db_path / "metadata.pkl"
        try:
            with open(str(metadata_path.absolute()), 'wb') as f:
                pickle.dump(self.documents, f)
            file_size = metadata_path.stat().st_size
            print(f"  ├─ ✅ 元数据已保存: {metadata_path.name} ({file_size:,} 字节)")
        except Exception as e:
            print(f"  ├─ ❌ 元数据保存失败: {e}")
        
        # 保存BM25索引数据（用于混合检索）
        if BM25_AVAILABLE and self.bm25 is not None:
            bm25_path = self.vector_db_path / "bm25_data.pkl"
            try:
                bm25_data = {
                    'tokenized_docs': self.tokenized_docs,
                    'bm25': self.bm25
                }
                with open(str(bm25_path.absolute()), 'wb') as f:
                    pickle.dump(bm25_data, f)
                file_size = bm25_path.stat().st_size
                print(f"  ├─ ✅ BM25索引已保存: {bm25_path.name} ({file_size:,} 字节)")
            except Exception as e:
                print(f"  ├─ ❌ BM25索引保存失败: {e}")
        
        # 保存配置信息和已处理文件记录
        config = {
            'model_name': self.model_name,
            'embedding_dim': self.embedding_dim,
            'total_documents': len(self.documents),
            'processed_files': self.processed_files
        }
        config_path = self.vector_db_path / "config.json"
        try:
            with open(str(config_path.absolute()), 'w', encoding='utf-8') as f:
                json.dump(config, f, indent=2, ensure_ascii=False)
            file_size = config_path.stat().st_size
            print(f"  └─ ✅ 配置信息已保存: {config_path.name} ({file_size:,} 字节)")
        except Exception as e:
            print(f"  └─ ❌ 配置保存失败: {e}")
        
        print("\n📋 验证保存结果...")
        all_success = True
        for filename in ['faiss.index', 'metadata.pkl', 'config.json']:
            filepath = self.vector_db_path / filename
            if filepath.exists():
                size = filepath.stat().st_size
                print(f"  ✅ {filename}: {size:,} 字节")
            else:
                print(f"  ❌ {filename}: 未找到")
                all_success = False
        
        if all_success:
            print("\n✅ 所有数据已成功保存并验证")
        else:
            print("\n⚠️  部分文件保存可能失败，请检查")
        
    def load_database(self):
        """从磁盘加载向量数据库（优化中文路径处理）"""
        print(f"\n📂 加载向量数据库从: {self.vector_db_path.absolute()}")
        
        index_path = self.vector_db_path / "faiss.index"
        metadata_path = self.vector_db_path / "metadata.pkl"
        config_path = self.vector_db_path / "config.json"
        
        # 检查文件是否存在
        if not all([index_path.exists(), metadata_path.exists(), config_path.exists()]):
            print("❌ 向量数据库文件不完整，请先构建数据库")
            return False
        
        # 加载FAISS索引（多方法处理中文路径问题）
        index_loaded = False
        
        # 方法1：使用Windows短路径名（8.3格式）
        if sys.platform == 'win32':
            try:
                abs_index_path = str(index_path.absolute())
                short_path = self.get_short_path_name(abs_index_path)
                
                if short_path != abs_index_path:
                    self.index = faiss.read_index(short_path)
                    print(f"  ├─ ✅ FAISS索引已加载（短路径）: {self.index.ntotal} 个向量")
                    index_loaded = True
            except Exception as e:
                print(f"  ├─ ⚠️  短路径方法失败: {e.__class__.__name__}")
        
        # 方法2：直接读取（如果短路径未成功）
        if not index_loaded:
            try:
                abs_index_path = str(index_path.absolute())
                self.index = faiss.read_index(abs_index_path)
                print(f"  ├─ ✅ FAISS索引已加载: {self.index.ntotal} 个向量")
                index_loaded = True
            except Exception as e:
                print(f"  ├─ ⚠️  直接加载失败: {e.__class__.__name__}")
        
        # 方法3：通过临时文件（最后的备用方法）
        if not index_loaded:
            try:
                print(f"  ├─ 使用临时文件备用方法...")
                import tempfile
                import shutil
                
                # 创建临时文件
                with tempfile.NamedTemporaryFile(mode='wb', delete=False, suffix='.index') as tmp:
                    temp_path = tmp.name
                
                # 复制到临时文件
                shutil.copy2(str(index_path.absolute()), temp_path)
                
                # 从临时文件加载
                self.index = faiss.read_index(temp_path)
                
                # 删除临时文件
                os.unlink(temp_path)
                print(f"  ├─ ✅ FAISS索引已加载（备用方法）: {self.index.ntotal} 个向量")
                index_loaded = True
            except Exception as e:
                print(f"  ├─ ❌ 备用方法失败: {e}")
                import traceback
                traceback.print_exc()
        
        if not index_loaded:
            print("❌ FAISS索引加载失败！所有方法均失败")
            return False
        
        # 加载元数据
        try:
            with open(str(metadata_path.absolute()), 'rb') as f:
                self.documents = pickle.load(f)
            print(f"  ├─ ✅ 元数据已加载: {len(self.documents)} 条记录")
        except Exception as e:
            print(f"  ├─ ❌ 元数据加载失败: {e}")
            return False
        
        # 加载BM25索引（用于混合检索）
        if BM25_AVAILABLE:
            bm25_path = self.vector_db_path / "bm25_data.pkl"
            if bm25_path.exists():
                try:
                    with open(str(bm25_path.absolute()), 'rb') as f:
                        bm25_data = pickle.load(f)
                    self.tokenized_docs = bm25_data['tokenized_docs']
                    self.bm25 = bm25_data['bm25']
                    print(f"  ├─ ✅ BM25索引已加载: {len(self.tokenized_docs)} 个文档")
                except Exception as e:
                    print(f"  ├─ ⚠️  BM25索引加载失败: {e}")
                    self.bm25 = None
            else:
                print(f"  ├─ ⚠️  未找到BM25索引文件，将只使用语义检索")
        
        # 加载配置和已处理文件记录
        try:
            with open(str(config_path.absolute()), 'r', encoding='utf-8') as f:
                config = json.load(f)
            print(f"  ├─ ✅ 配置已加载: 向量维度 {config['embedding_dim']}")
            # 加载已处理文件记录
            self.processed_files = config.get('processed_files', {})
            if self.processed_files:
                print(f"  └─ ✅ 已处理文件记录: {len(self.processed_files)} 个文件")
            else:
                print(f"  └─ ℹ️  无已处理文件记录（旧版数据库）")
        except Exception as e:
            print(f"  └─ ⚠️  配置加载失败: {e}")
            self.processed_files = {}
        
        print("✅ 向量数据库加载完成\n")
        return True
    
    def search(self, query: str, top_k: int = 5) -> List[Dict[str, Any]]:
        """
        语义搜索
        
        Args:
            query: 查询文本
            top_k: 返回前k个最相似的结果
            
        Returns:
            搜索结果列表，包含文档信息和相似度分数
        """
        if self.index is None or not self.documents:
            print("❌ 向量数据库未初始化，请先构建或加载数据库")
            return []
        
        # 将查询文本向量化
        query_vector = self.embed_texts([query])
        
        # FAISS搜索
        distances, indices = self.index.search(query_vector, top_k)
        
        # 构建结果
        results = []
        for idx, (dist, doc_idx) in enumerate(zip(distances[0], indices[0]), 1):
            if doc_idx < len(self.documents):
                doc = self.documents[doc_idx].copy()
                # L2距离转换为相似度分数 (越小越相似)
                doc['similarity_score'] = float(1 / (1 + dist))
                doc['distance'] = float(dist)
                doc['rank'] = idx
                results.append(doc)
        
        return results
    
    def hybrid_search(self, 
                     query: str, 
                     top_k: int = 5,
                     semantic_weight: float = 0.7,
                     keyword_weight: float = 0.3,
                     rerank: bool = True) -> List[Dict[str, Any]]:
        """
        混合检索：语义检索 + BM25关键词检索 + 可选Rerank
        
        Args:
            query: 查询文本
            top_k: 返回前k个最相似的结果
            semantic_weight: 语义检索权重（默认0.7）
            keyword_weight: 关键词检索权重（默认0.3）
            rerank: 是否使用LLM进行重排序
            
        Returns:
            混合检索结果列表
        """
        if self.index is None or not self.documents:
            print("❌ 向量数据库未初始化，请先构建或加载数据库")
            return []
        
        # 1. 语义检索（FAISS）
        semantic_results = self.search(query, top_k=top_k * 2)  # 多检索一些候选
        
        # 2. BM25关键词检索（如果可用）
        if BM25_AVAILABLE and self.bm25 is not None:
            # 对查询分词
            query_tokens = TextChunker.tokenize(query)
            
            # BM25检索
            bm25_scores = self.bm25.get_scores(query_tokens)
            
            # 获取top_k*2个BM25结果
            top_indices = np.argsort(bm25_scores)[::-1][:top_k * 2]
            
            # 构建BM25结果
            bm25_results = []
            for idx in top_indices:
                if idx < len(self.documents):
                    doc = self.documents[idx].copy()
                    doc['bm25_score'] = float(bm25_scores[idx])
                    bm25_results.append(doc)
            
            # 3. 混合打分
            # 创建文档ID到分数的映射
            semantic_scores = {}
            for doc in semantic_results:
                doc_id = (doc['filename'], doc['chunk_id'])
                semantic_scores[doc_id] = doc['similarity_score']
            
            bm25_score_dict = {}
            for doc in bm25_results:
                doc_id = (doc['filename'], doc['chunk_id'])
                bm25_score_dict[doc_id] = doc['bm25_score']
            
            # 归一化BM25分数到[0,1]
            if bm25_score_dict:
                max_bm25 = max(bm25_score_dict.values())
                min_bm25 = min(bm25_score_dict.values())
                if max_bm25 > min_bm25:
                    bm25_score_dict = {
                        k: (v - min_bm25) / (max_bm25 - min_bm25) 
                        for k, v in bm25_score_dict.items()
                    }
            
            # 合并所有文档ID
            all_doc_ids = set(semantic_scores.keys()) | set(bm25_score_dict.keys())
            
            # 计算混合分数
            hybrid_results = []
            for doc_id in all_doc_ids:
                sem_score = semantic_scores.get(doc_id, 0.0)
                bm25_score = bm25_score_dict.get(doc_id, 0.0)
                
                # 混合分数
                hybrid_score = (semantic_weight * sem_score + 
                              keyword_weight * bm25_score)
                
                # 找到对应的文档
                doc = None
                for d in semantic_results + bm25_results:
                    if (d['filename'], d['chunk_id']) == doc_id:
                        doc = d.copy()
                        break
                
                if doc:
                    doc['hybrid_score'] = hybrid_score
                    doc['semantic_score'] = sem_score
                    doc['bm25_score'] = bm25_score
                    hybrid_results.append(doc)
            
            # 按混合分数排序
            hybrid_results.sort(key=lambda x: x['hybrid_score'], reverse=True)
            hybrid_results = hybrid_results[:top_k * 2]  # 保留2倍候选用于rerank
            
        else:
            # 如果BM25不可用，直接使用语义检索结果
            hybrid_results = semantic_results
            for doc in hybrid_results:
                doc['hybrid_score'] = doc['similarity_score']
        
        # 4. Rerank（使用LLM重排序）
        if rerank and len(hybrid_results) > top_k:
            hybrid_results = self._rerank_with_llm(query, hybrid_results, top_k)
        else:
            hybrid_results = hybrid_results[:top_k]
        
        # 更新最终排名
        for idx, doc in enumerate(hybrid_results, 1):
            doc['rank'] = idx
            # 使用hybrid_score作为最终的similarity_score
            doc['similarity_score'] = doc.get('hybrid_score', doc.get('similarity_score', 0.0))
        
        return hybrid_results
    
    def _rerank_with_llm(self, query: str, candidates: List[Dict], top_k: int) -> List[Dict]:
        """
        使用LLM对候选结果进行重排序
        
        Args:
            query: 用户查询
            candidates: 候选文档列表
            top_k: 返回前k个结果
            
        Returns:
            重排序后的结果
        """
        try:
            from dashscope import Generation
            
            # 构建重排序prompt
            prompt = f"""请评估以下文档片段与用户问题的相关性，给每个文档打分（0-10分）。

用户问题：{query}

文档片段：
"""
            for idx, doc in enumerate(candidates[:10], 1):  # 最多评估10个
                text_preview = doc['text'][:200]
                prompt += f"\n[{idx}] {text_preview}...\n"
            
            prompt += "\n请以JSON格式输出每个文档的相关性分数，例如：{\"1\": 8.5, \"2\": 6.0, ...}"
            
            # 调用LLM
            response = Generation.call(
                model="qwen-turbo",
                prompt=prompt,
                temperature=0.1,
                max_tokens=500
            )
            
            if response.status_code == 200:
                import re
                import json
                
                # 提取JSON
                text = response.output['text']
                json_match = re.search(r'\{[^}]+\}', text)
                
                if json_match:
                    scores = json.loads(json_match.group())
                    
                    # 应用LLM分数
                    for idx_str, score in scores.items():
                        idx = int(idx_str) - 1
                        if 0 <= idx < len(candidates):
                            candidates[idx]['llm_rerank_score'] = float(score)
                    
                    # 按LLM分数重排序
                    candidates.sort(
                        key=lambda x: x.get('llm_rerank_score', 0), 
                        reverse=True
                    )
        
        except Exception as e:
            print(f"  ⚠️  LLM Rerank失败: {e}，使用原始排序")
        
        return candidates[:top_k]
    
    def add_documents(self, 
                     file_paths: List[str],
                     chunk_size: int = 500,
                     overlap: int = 50):
        """
        增量添加新文档到现有向量库
        
        Args:
            file_paths: 要添加的文件路径列表
            chunk_size: 文本块大小
            overlap: 块之间重叠大小
        """
        print("\n" + "="*60)
        print("📚 增量添加文档到向量库")
        print("="*60)
        
        if self.index is None:
            print("⚠️  向量库未初始化，将创建新库")
            self.index = faiss.IndexFlatL2(self.embedding_dim)
            self.documents = []
        
        all_chunks = []
        all_metadata = []
        processed_count = 0
        skipped_count = 0
        
        for file_path_str in file_paths:
            file_path = Path(file_path_str)
            
            if not file_path.exists():
                print(f"⚠️  文件不存在: {file_path.name}")
                continue
            
            # 检查是否已处理
            if self.is_file_processed(file_path):
                print(f"⏭️  跳过已处理文件: {file_path.name}")
                skipped_count += 1
                continue
            
            print(f"\n📄 处理文档: {file_path.name}")
            
            # 解析文档
            print("  ├─ 解析文档内容...")
            text = DocumentParser.parse_document(str(file_path))
            
            if not text.strip():
                print("  └─ ⚠️  文档内容为空，跳过")
                continue
            
            print(f"  ├─ 提取文本: {len(text)} 字符")
            
            # 文本分块
            print(f"  ├─ 文本分块...")
            chunks = TextChunker.chunk_text(text, chunk_size, overlap)
            print(f"  └─ ✅ 生成 {len(chunks)} 个文本块")
            
            # 保存块和元数据
            for chunk_idx, chunk in enumerate(chunks):
                all_chunks.append(chunk)
                all_metadata.append({
                    'filename': file_path.name,
                    'chunk_id': chunk_idx,
                    'text': chunk
                })
            
            # 记录文件为已处理
            self.processed_files[file_path.name] = {
                'hash': self.get_file_hash(str(file_path)),
                'mtime': file_path.stat().st_mtime,
                'chunk_count': len(chunks)
            }
            processed_count += 1
        
        if not all_chunks:
            print(f"\n⏭️  没有新文件需要处理")
            print(f"  • 已处理: {skipped_count} 个")
            return
        
        print(f"\n{'='*60}")
        print(f"📊 文本处理完成")
        print(f"  • 新增文档: {processed_count} 个")
        print(f"  • 跳过文档: {skipped_count} 个")
        print(f"  • 新增文本块: {len(all_chunks)} 个")
        print(f"{'='*60}")
        
        # 向量化
        print("\n🔄 使用DashScope API向量化文本块...")
        embeddings = self.embed_texts(all_chunks, batch_size=25)
        print(f"✅ 向量化完成，向量形状: {embeddings.shape}")
        
        # 添加到FAISS索引
        print("\n🔨 更新FAISS索引...")
        self.index.add(embeddings.astype('float32'))
        self.documents.extend(all_metadata)
        print(f"✅ FAISS索引更新完成，现共 {self.index.ntotal} 个向量")
        
        # 保存向量库
        self.save_database()
        
        print(f"\n{'='*60}")
        print("🎉 文档添加完成！")
        print(f"{'='*60}")
    
    def list_processed_files(self):
        """列出所有已处理的文件"""
        if not self.processed_files:
            print("📋 暂无已处理文件记录")
            return
        
        print(f"\n📋 已处理文件列表 ({len(self.processed_files)} 个)")
        print("="*60)
        for filename, info in self.processed_files.items():
            chunk_count = info.get('chunk_count', 0)
            print(f"  📄 {filename}")
            print(f"     └─ 文本块数: {chunk_count}")
        print("="*60)


def main():
    """主函数：构建向量数据库"""
    
    print("="*60)
    print("🚀 嵌入式RAG向量数据库构建系统")
    print("="*60)
    print("📌 使用DashScope Text Embedding API")
    print("📌 API Key从环境变量 DASHSCOPE_API_KEY 读取")
    print("="*60 + "\n")
    
    # 初始化数据库
    db = EmbeddingVectorDatabase(
        model_name="text-embedding-v2",
        vector_db_path="./embedding向量库"
    )
    
    # 构建向量数据库
    db.build_database(
        data_folder=r"E:\AI_Code\嵌入式硬件RAG问答Agent\data",
        chunk_size=500,    # 文本块大小：适合技术文档
        overlap=50          # 重叠大小：保证上下文连贯性
    )
    
    print("\n" + "="*60)
    print("🎯 测试搜索功能")
    print("="*60)
    
    # 测试搜索
    test_queries = [
        "激光雷达的测距范围是多少",
        "ADS6311的主要特性",
        "Hawk模组的规格参数"
    ]
    
    for query in test_queries:
        print(f"\n🔍 查询: {query}")
        results = db.search(query, top_k=3)
        
        for result in results:
            print(f"\n  📄 文件: {result['filename']}")
            print(f"  📊 相似度: {result['similarity_score']:.4f}")
            print(f"  📝 内容预览: {result['text'][:150]}...")
            print("  " + "-"*56)


if __name__ == "__main__":
    main()

